from pptx import Presentation
import os
import re

# Ruta al archivo PowerPoint
pptx_path = os.path.expanduser(r'C:\Github\Ampa\documentos\Cartones.pptx')

print(f"Buscando archivo en: {pptx_path}")

if not os.path.exists(pptx_path):
    print(f"❌ ERROR: No se encontró el archivo {pptx_path}")
    exit(1)

# Abrir la presentación
print("📂 Abriendo presentación...")
prs = Presentation(pptx_path)

# Extraer todo el texto de cada diapositiva
slides_data = []
for slide in prs.slides:
    texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]
    slides_data.append(texts)

print(f"✅ Se encontraron {len(slides_data)} diapositivas")

# Generar el markdown
md_lines = ['# Cartones de Bingo Musical', '']
card_index = 1

# Patrón para detectar canciones: ID + espacio + artista - canción
# Ejemplo: "4 Coyote Dax - No Rompas Más"
song_pattern = re.compile(r'^\d+\s+.+\s+-\s+.+')

for i, texts in enumerate(slides_data, start=1):
    md_lines.append(f'# Diapositiva {i}')
    md_lines.append('')
    
    # Filtrar solo las líneas que son canciones válidas
    songs = [t for t in texts if song_pattern.match(t)]
    
    # IMPORTANTE: Eliminar duplicados manteniendo el orden
    unique_songs = []
    seen = set()
    for song in songs:
        if song not in seen:
            unique_songs.append(song)
            seen.add(song)
    
    songs = unique_songs
    
    print(f"  Diapositiva {i}: {len(songs)} canciones únicas detectadas")
    
    # Verificar que sea divisible por 12 (debería haber 3 cartones por diapositiva = 36 canciones)
    n = len(songs)
    if n != 36:
        print(f"    ⚠️  ADVERTENCIA: {n} canciones (se esperaban 36 para 3 cartones)")
    
    # Agrupar de 12 en 12
    num_cards = 3  # Siempre 3 cartones por diapositiva
    
    for card_num in range(num_cards):
        start_idx = card_num * 12
        end_idx = min(start_idx + 12, n)
        card_songs = songs[start_idx:end_idx]
        
        md_lines.append(f'## Cartón {card_index}')
        md_lines.append('')
        
        # Verificar que tenga exactamente 12 canciones
        if len(card_songs) != 12:
            print(f"    ⚠️  ERROR: Cartón {card_index} tiene {len(card_songs)} canciones en lugar de 12")
        
        for idx, entry in enumerate(card_songs, start=1):
            md_lines.append(f'{idx} - {entry}')
        
        md_lines.append('')
        md_lines.append('---')
        card_index += 1
    
    md_lines.append('')

# Guardar en archivo en el directorio padre del script
script_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(script_dir)
output_file = os.path.join(parent_dir, 'cartones-extraidos.md')

md_text = '\n'.join(md_lines)
with open(output_file, 'w', encoding='utf-8') as f:
    f.write(md_text)

print(f'\n✅ Procesadas {len(slides_data)} diapositivas')
print(f'✅ Extraídos {card_index - 1} cartones')
print(f'✅ Archivo guardado: {output_file}')
print(f'\n💡 Revisa las advertencias arriba para identificar diapositivas con problemas.')
